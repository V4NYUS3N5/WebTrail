"""
生成 WebTrail 项目讲解 PPT（约 18 页）。
蓝白主题 · 依赖: pip install python-pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# ── 蓝白主题配色 ─────────────────────────────────────
CLR_BG      = RGBColor(0xFF, 0xFF, 0xFF)  # 白色背景
CLR_BLUE    = RGBColor(0x1A, 0x73, 0xE8)  # 主蓝色
CLR_DBLUE   = RGBColor(0x0D, 0x47, 0xA1)  # 深蓝（标题文字）
CLR_LBLUE   = RGBColor(0xE3, 0xF2, 0xFD)  # 极浅蓝（卡片底色）
CLR_WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
CLR_TEXT    = RGBColor(0x33, 0x33, 0x33)  # 正文深灰
CLR_MUTED   = RGBColor(0x66, 0x66, 0x66)  # 次要文字
CLR_CODE_BG = RGBColor(0xF0, 0xF4, 0xF8)  # 代码块背景
CLR_CODE_TX = RGBColor(0x1A, 0x2A, 0x3A)  # 代码文字
CLR_TITLE_BAR_BG = RGBColor(0x15, 0x65, 0xC0)  # 标题栏深蓝
CLR_GREEN   = RGBColor(0x2E, 0x7D, 0x32)
CLR_ORANGE  = RGBColor(0xE6, 0x51, 0x00)
CLR_RED     = RGBColor(0xC6, 0x28, 0x28)
CLR_BORDER  = RGBColor(0xBB, 0xDE, 0xFB)  # 卡片边框浅蓝
CLR_CARD_BG = RGBColor(0xF8, 0xFB, 0xFF)  # 卡片微蓝底
CLR_COVER_BG = RGBColor(0x15, 0x65, 0xC0) # 封面/尾页蓝色底色

FONT_TITLE = "Microsoft YaHei"
FONT_BODY  = "Microsoft YaHei"
FONT_MONO  = "Consolas"


# ═══════════════════════════════════════
# 绘图工具函数
# ═══════════════════════════════════════

def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_title_bar(slide, text: str, subtitle: str = ""):
    """顶部标题条。"""
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                 Inches(0), Inches(0), Inches(13.33), Inches(1.1))
    bar.fill.solid()
    bar.fill.fore_color.rgb = CLR_TITLE_BAR_BG
    bar.line.fill.background()

    tf = bar.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = CLR_WHITE
    p.font.name = FONT_TITLE
    p.alignment = PP_ALIGN.LEFT
    tf.margin_left = Inches(0.6)
    tf.margin_top = Inches(0.18)

    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                  Inches(0), Inches(1.1), Inches(13.33), Pt(4))
    line.fill.solid()
    line.fill.fore_color.rgb = CLR_BLUE
    line.line.fill.background()

    if subtitle:
        tb = slide.shapes.add_textbox(Inches(0.7), Inches(1.3), Inches(11.9), Inches(0.5))
        tf2 = tb.text_frame
        p2 = tf2.paragraphs[0]
        p2.text = subtitle
        p2.font.size = Pt(13)
        p2.font.color.rgb = CLR_MUTED
        p2.font.name = FONT_BODY


def add_body_text(slide, left, top, width, height, lines, font_size=14,
                  color=None, line_spacing=1.5):
    if color is None:
        color = CLR_TEXT
    tb = slide.shapes.add_textbox(Inches(left), Inches(top),
                                  Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = FONT_BODY
        p.space_after = Pt(6)
        p.line_spacing = Pt(font_size * line_spacing)
    return tb


def add_card(slide, left, top, width, height, title: str, body_lines):
    """带标题的白色卡片 + 浅蓝底。"""
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                  Inches(left), Inches(top),
                                  Inches(width), Inches(height))
    card.fill.solid()
    card.fill.fore_color.rgb = CLR_CARD_BG
    card.line.color.rgb = CLR_BORDER
    card.line.width = Pt(1.5)

    tf = card.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.2)
    tf.margin_right = Inches(0.2)
    tf.margin_top = Inches(0.15)

    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = CLR_DBLUE
    p.font.name = FONT_TITLE
    p.space_after = Pt(6)

    for line in body_lines:
        p2 = tf.add_paragraph()
        p2.text = line
        p2.font.size = Pt(11)
        p2.font.color.rgb = CLR_TEXT
        p2.font.name = FONT_BODY
        p2.space_after = Pt(3)

    return card


def add_code_block(slide, left, top, width, height, code: str):
    """代码展示块（浅灰蓝底）。"""
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                 Inches(left), Inches(top),
                                 Inches(width), Inches(height))
    box.fill.solid()
    box.fill.fore_color.rgb = CLR_CODE_BG
    box.line.color.rgb = CLR_BORDER
    box.line.width = Pt(1)

    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.25)
    tf.margin_top = Inches(0.15)

    for i, line in enumerate(code.strip().split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(11)
        p.font.color.rgb = CLR_CODE_TX
        p.font.name = FONT_MONO
        p.space_after = Pt(2)
    return box


def add_page_number(slide, num, total, dark=False):
    color = CLR_WHITE if dark else CLR_MUTED
    tb = slide.shapes.add_textbox(Inches(12), Inches(7.2), Inches(1.1), Inches(0.3))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = f"{num}/{total}"
    p.font.size = Pt(10)
    p.font.color.rgb = color
    p.font.name = FONT_BODY
    p.alignment = PP_ALIGN.RIGHT


def cover_footer(slide, lines):
    add_body_text(slide, 1, 6.0, 11, 1.3, lines, font_size=14, color=CLR_WHITE)


# ═══════════════════════════════════════
# 主生成
# ═══════════════════════════════════════

def build():
    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)
    TOTAL = 18
    BLANK = prs.slide_layouts[6]

    def new_slide(num):
        slide = prs.slides.add_slide(BLANK)
        set_slide_bg(slide, CLR_BG)
        add_page_number(slide, num, TOTAL)
        return slide

    # ── 1. 封面 ──
    s = prs.slides.add_slide(BLANK)
    set_slide_bg(s, CLR_COVER_BG)
    add_body_text(s, 1, 2.0, 11, 1.2, ["WebTrail"],
                  font_size=60, color=CLR_WHITE)
    add_body_text(s, 1, 3.2, 11, 0.8, ["浏览器数字取证与用户画像工具"],
                  font_size=28, color=CLR_LBLUE)
    add_body_text(s, 1, 4.4, 11, 0.6,
                  ["纯 Python 标准库 · 零第三方依赖 · 开箱即用"],
                  font_size=16, color=CLR_WHITE)
    add_body_text(s, 1, 5.4, 11, 0.5,
                  ["命令行 + 图形界面 · 跨平台 Windows / macOS / Linux"],
                  font_size=14, color=CLR_LBLUE)
    add_page_number(s, 1, TOTAL, dark=True)

    # ── 2. 项目概述 ──
    s = new_slide(2)
    add_title_bar(s, "项目概述", "WebTrail 是什么？")
    add_body_text(s, 0.7, 1.7, 11.9, 5, [
        "WebTrail 是一款面向数字取证的开源工具，从本地浏览器中提取完整的浏览痕迹，",
        "并通过用户行为画像技术，将碎片化数据转化为结构化的调查情报。",
        "",
        "  ▸ 目标用户：  安全分析师、取证调查员、IT 审计人员",
        "  ▸ 核心价值：  提取 + 分析 一体化，从数据采集到报告输出一条链路",
        "  ▸ 技术特点：  纯 Python 标准库实现，无需安装任何第三方包",
        "  ▸ 运行方式：  命令行（CLI）和图形界面（GUI）双模式",
        "  ▸ 跨平台：    Windows / macOS / Linux 均可运行",
    ], font_size=15)

    # ── 3. 背景与动机 ──
    s = new_slide(3)
    add_title_bar(s, "背景与动机", "为什么需要这样一款工具？")
    add_card(s, 0.7, 1.7, 3.7, 2.2, "痛点 1: 工具碎片化", [
        "  商业取证工具（EnCase / FTK）",
        "  价格高昂、部署复杂、捆绑安装",
        "  需要便携、绿色、开箱即用的方案",
    ])
    add_card(s, 4.8, 1.7, 3.7, 2.2, "痛点 2: 信息过载", [
        "  传统工具只做「提取」这一步",
        "  调查者面对数千条原始记录",
        "  信息密度低，人工研判耗时",
    ])
    add_card(s, 8.9, 1.7, 3.7, 2.2, "痛点 3: 缺乏画像能力", [
        "  需要从痕迹中自动提炼行为模式",
        "  活跃时段、兴趣偏好、风险行为",
        "  将数据筛选转化为可操作情报",
    ])
    add_body_text(s, 0.7, 4.3, 11.9, 2.5, [
        "",
        "  WebTrail 的答案：一键完成 证据固定 → 数据提取 → 画像分析 → 报告输出",
        "  将取证工作流压缩为一个命令：  python main.py",
    ], font_size=14, color=CLR_DBLUE)

    # ── 4. 核心功能总览 ──
    s = new_slide(4)
    add_title_bar(s, "核心功能总览")
    add_card(s, 0.7, 1.7, 2.8, 2.5, "痕迹提取", [
        "  浏览历史记录",
        "  Cookie 信息",
        "  下载记录",
        "  书签数据",
        "  登录凭证（密码）",
        "  缓存 / HSTS 元信息",
    ])
    add_card(s, 3.8, 1.7, 2.8, 2.5, "浏览器支持", [
        "  Google Chrome",
        "  Microsoft Edge",
        "  Mozilla Firefox",
        "  （Chromium 内核 Edge",
        "   通过继承复用代码）",
    ])
    add_card(s, 6.9, 1.7, 2.8, 2.5, "取证保障", [
        "  只读安全访问",
        "  SHA-256 链式保管",
        "  记录级完整性校验",
        "  时间戳保真转换",
        "  临时文件自动清理",
    ])
    add_card(s, 10.0, 1.7, 2.8, 2.5, "画像分析", [
        "  24h 活跃热力图",
        "  TOP 域名排名",
        "  URL 语义分类（8类）",
        "  风险指标检测",
        "  多浏览器对比",
    ])

    # ── 5. 技术架构总览 ──
    s = new_slide(5)
    add_title_bar(s, "技术架构总览", "5 层模块化设计 · 15 个源文件 · 约 1800 行代码")
    add_body_text(s, 0.7, 1.7, 11.9, 0.5,
                  ["分层架构，每层职责单一："], font_size=15, color=CLR_DBLUE)

    layers = [
        ("入口层", ["main.py (87行)", "解析 CLI 参数", "调度 GUI 或管道"]),
        ("管道层", ["pipeline.py (119行)", "三阶段流程编排", "哈希 → 提取 → 画像"]),
        ("适配器层", ["browsers/", "Chrome / Edge / Firefox", "ABC 抽象基类 + 继承"]),
        ("分析层", ["core/", "profiler.py 七维度画像", "timeline.py 时间线"]),
        ("输出层", ["output/", "JSON 取证报告", "CSV 时间线表格"]),
    ]
    for i, (name, desc) in enumerate(layers):
        add_card(s, 0.5 + i * 2.55, 2.3, 2.35, 2.0, name, desc)

    add_body_text(s, 0.7, 4.7, 11.9, 2.5, [
        "设计原则：  关注点分离  ·  最少依赖  ·  错误隔离  ·  职责单一  ·  继承合理",
        "",
        "main.py → pipeline.py → browsers/* → profiler.py → output/writer.py",
    ], font_size=14, color=CLR_DBLUE)

    # ── 6. 三阶段取证管道 ──
    s = new_slide(6)
    add_title_bar(s, "三阶段取证管道", "核心编排逻辑 — pipeline.py")

    stages = [
        ("阶段 1", "证据哈希", "Chain of Custody",
         ["▸ 遍历浏览器目录下所有 .sqlite 和 .json",
          "▸ 逐文件计算 SHA-256 校验和",
          "▸ 记录到报告元数据中",
          "▸ 示例：189 个证据文件全部哈希"],
         CLR_GREEN),
        ("阶段 2", "痕迹提取", "Artifact Extraction",
         ["▸ 每个浏览器适配器并行运行",
          "▸ 统一数据结构 ArtifactRecord",
          "▸ 错误隔离：Chrome 失败不影响 Firefox",
          "▸ 记录级校验和保证单条完整性"],
         CLR_BLUE),
        ("阶段 3", "用户画像", "Behavior Profiling",
         ["▸ 基于全部痕迹的七维度分析",
          "▸ 活跃热力图 / TOP 域名 / URL 分类",
          "▸ 风险指标自动检测",
          "▸ 生成 JSON + CSV 结构化报告"],
         CLR_ORANGE),
    ]
    for i, (tag, cn, en, lines, clr) in enumerate(stages):
        y = 1.7 + i * 2.0
        tag_box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                     Inches(0.7), Inches(y), Inches(1.5), Inches(0.45))
        tag_box.fill.solid()
        tag_box.fill.fore_color.rgb = clr
        tag_box.line.fill.background()

        tf = tag_box.text_frame
        p = tf.paragraphs[0]
        p.text = tag
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = CLR_WHITE
        p.font.name = FONT_TITLE
        p.alignment = PP_ALIGN.CENTER

        add_body_text(s, 2.5, y - 0.05, 10, 0.5,
                      [f"{cn}（{en}）"], font_size=17, color=CLR_DBLUE)
        add_body_text(s, 2.7, y + 0.55, 9.5, 1.3, lines, font_size=13)

    # ── 7. 阶段 1：证据哈希 ──
    s = new_slide(7)
    add_title_bar(s, "阶段 1：证据文件哈希", "链式保管（Chain of Custody）")
    add_body_text(s, 0.7, 1.7, 11.9, 1.0, [
        "在提取任何数据之前，对全部证据文件计算 SHA-256。确保数据从未被篡改。",
    ], font_size=15, color=CLR_TEXT)
    add_code_block(s, 0.7, 2.5, 6.0, 3.5, """\
def collect_evidence_hashes(extractors) -> dict:
    \"\"\"阶段1：收集证据文件哈希\"\"\"
    all_hashes = {}
    for ext in extractors:
        for name, path in ext.detect_profiles():
            for f in path.rglob("*"):
                if f.suffix in (".sqlite", ".json"):
                    label = f"{ext.browser}/{name}/{f.name}"
                    all_hashes[label] = file_hash(f)
    return all_hashes""")
    add_card(s, 7.2, 2.5, 5.5, 3.5, "哈希链验证流程", [
        "",
        "  ┌───────────────┐",
        "  │ 证据文件        │",
        "  │ .sqlite  .json  │",
        "  └───────┬───────┘",
        "          │ SHA-256",
        "          ▼",
        "  ┌───────────────┐",
        "  │ report.json    │ ← 元数据记录全部哈希",
        "  │ metadata:       │",
        "  │  {file_hash}    │",
        "  └───────┬───────┘",
        "          │ SHA-256",
        "          ▼",
        "  ┌──────────────────┐",
        "  │ report_checksums  │ ← 报告自校验",
        "  └──────────────────┘",
    ])

    # ── 8. 阶段 2：痕迹提取 ──
    s = new_slide(8)
    add_title_bar(s, "阶段 2：浏览器痕迹提取", "统一数据结构 · 错误隔离")
    add_code_block(s, 0.7, 1.7, 6.0, 2.2, """\
@dataclass
class ArtifactRecord:
    artifact_type: str     # history|cookie|download|bookmark|login|cache
    browser: str           # Chrome|Edge|Firefox
    timestamp:  str|None   # ISO 8601
    profile:    str        # 用户配置名
    source_file: str       # 证据文件路径
    data:       dict       # 具体载荷
    extraction_time: str   # 提取时刻
    checksum:   str        # 本条记录 SHA-256""")
    add_card(s, 7.2, 1.7, 5.5, 2.2, "六大提取类型", [
        "  History   → SQLite (urls / moz_places)",
        "  Cookie    → SQLite (cookies / moz_cookies)",
        "  Download  → SQLite (downloads / moz_annos)",
        "  Bookmark  → JSON 递归解析",
        "  Login     → SQLite / logins.json",
        "  Cache     → 目录元信息 + HSTS",
    ])
    add_body_text(s, 0.7, 4.3, 11.9, 2.7, [
        "提取保障机制：",
        "  ▸ SQLite 连接全部 mode=ro + PRAGMA query_only=ON",
        "  ▸ 浏览器运行导致锁定 → safe_connect_with_fallback() 自动复制到临时文件，关闭即清理",
        "  ▸ 单条提取失败不中断其他提取器（try/except per extractor）",
        "  ▸ ExtractionResult 包含 records + errors 双通道",
    ], font_size=14)

    # ── 9. 阶段 3：用户画像 ──
    s = new_slide(9)
    add_title_bar(s, "阶段 3：用户行为画像", "七维度综合分析 — profiler.py")

    dims = [
        ("overview", "总览", "记录总数、时间范围、\n浏览器分布、类型分布"),
        ("activity_heat", "活跃热力图", "24小时分布 + \n近30天逐日统计"),
        ("top_domains", "TOP 域名", "TOP 20 访问域名\n及访问频次排名"),
        ("top_categories", "URL 分类", "8类语义分类：\n社交/搜索/视频等"),
        ("behavior_insights", "行为洞察", "活跃高峰、夜间占比、\n登录/密码保存数量"),
        ("browser_usage", "浏览器使用", "各浏览器按痕迹类型\n细分统计"),
        ("risk_indicators", "风险指标", "隐私模式推断、历史清除\n检测、可疑域名匹配"),
    ]
    for i, (key, cn, desc) in enumerate(dims):
        x = 0.5 + (i % 4) * 3.15
        y = 1.7 + (i // 4) * 2.6
        add_card(s, x, y, 2.9, 2.3, f"{key}", [
            f"中文: {cn}",
            f"内容: {desc}",
        ])

    # ── 10. 浏览器适配器设计 ──
    s = new_slide(10)
    add_title_bar(s, "浏览器适配器设计", "ABC 抽象基类 + 继承复用")
    add_code_block(s, 0.7, 1.7, 6.0, 5.0, """\
class BaseExtractor(ABC):
    \"\"\"所有浏览器提取器的抽象基类\"\"\"
    browser: str
    base_path: Path

    def run(self) -> ExtractionResult:
        profiles = self.detect_profiles()
        for name, path in profiles:
            for extract in self._get_methods():
                try:
                    records += extract(path)
                except Exception as e:
                    errors.append(str(e))
        return ExtractionResult(...)

    @abstractmethod
    def detect_profiles(self): ...
    @abstractmethod
    def extract_history(self, path): ...
    @abstractmethod
    def extract_cookies(self, path): ...
    # ... 共 8 个抽象方法""")
    add_card(s, 7.2, 1.7, 5.5, 5.0, "适配器一览", [
        "  Chrome (212行)",
        "  ├─ detect_profiles: 扫描 Default / Profile N",
        "  ├─ 时间戳: 1601 epoch 微秒（WebKit）",
        "  └─ SQLite → History / Cookies / Login Data",
        "",
        "  Edge (18行 — 继承Chrome)",
        "  └─ 仅覆写 __init__ 设置 browser=\"Edge\"",
        "",
        "  Firefox (198行)",
        "  ├─ detect_profiles: 解析 profiles.ini",
        "  ├─ 时间戳: 1970 epoch 微秒（PRTime）",
        "  └─ SQLite → places.sqlite / cookies.sqlite",
    ])

    # ── 11. Edge 继承设计 ──
    s = new_slide(11)
    add_title_bar(s, 'Edge 继承 Chrome — 18 行复用 212 行逻辑', "OOP 继承最佳实践")
    add_code_block(s, 0.7, 1.7, 6.0, 2.5, """\
from .chrome import ChromeExtractor
from config import EDGE_BASE

class EdgeExtractor(ChromeExtractor):
    \"\"\"继承 Chrome 全部提取逻辑，仅覆写浏览器名。\"\"\"
    def __init__(self, base_path: Path | None = None):
        super().__init__(base_path or EDGE_BASE)
        self.browser = "Edge"    # 只有这行不同！""")
    add_body_text(s, 0.7, 4.5, 5.5, 2.5, [
        "为什么可以这样做？",
        "  ▸ Edge 使用 Chromium 内核，数据库格式与 Chrome 完全一致",
        "  ▸ History / Cookies / Login Data 表结构相同",
        "  ▸ 时间戳格式相同（WebKit 微秒）",
        "  ▸ 书签 JSON 结构相同",
        "",
        "  这是「组合优于继承」的反例——",
        "  当两个类有 95% 相同行为时，继承是最优解。",
    ], font_size=13)
    add_card(s, 7.2, 1.7, 5.5, 4.5, "不用继承 vs 用继承", [
        "  如果不用继承：",
        "  - 需复制粘贴 Chrome 212 行代码到 Edge",
        "  - SQL 查询、JSON 解析全部重复",
        "  - 修 bug 要改两个文件",
        "",
        "  使用继承后：",
        "  - Edge 只需 18 行",
        "  - 自动获得 Chrome 全部能力",
        "  - 修一个地方，两个浏览器同时受益",
        "  - 语义清晰：Edge 就是 Chromium 内核浏览器",
    ])

    # ── 12. 取证科学性 ──
    s = new_slide(12)
    add_title_bar(s, "取证科学性保障", "司法级严谨设计")
    items = [
        ("只读防污染",
         ["所有 SQLite 连接 mode=ro + PRAGMA query_only=ON",
          "浏览器运行时数据库被锁 → 自动 copy_to_temp 回退读取",
          "临时文件读取完毕后自动 unlink 清理，不留痕迹"]),
        ("链式保管",
         ["提取前对全部 .sqlite / .json 证据文件做 SHA-256",
          "哈希记录于 report.json 元数据",
          "报告本身也有 SHA-256 自校验 → 完整哈希链"]),
        ("记录级完整性",
         ["每条 ArtifactRecord 包含 checksum 字段",
          "serialize → JSON → SHA-256",
          "可逐条独立验证，防止选择性篡改"]),
        ("时间戳保真",
         ["Chrome: 1601-01-01 微秒数 → ISO 8601",
          "Firefox: 1970-01-01 微秒数 → ISO 8601",
          "保留原始时区信息 +00:00，不做截断",
          "画像时自动转本地时区显示"]),
    ]
    for i, (title, lines) in enumerate(items):
        x = 0.5 + (i % 2) * 6.3
        y = 1.7 + (i // 2) * 2.7
        add_card(s, x, y, 6.0, 2.4, title, lines)

    # ── 13. 风险指标检测 ──
    s = new_slide(13)
    add_title_bar(s, "风险指标检测机制", "三类自动风险检测")
    add_code_block(s, 0.7, 1.7, 6.5, 5.0, """\
def _risk_indicators(records) -> dict:
    indicators = {}

    # 1. 隐私模式痕迹推断
    indicators["private_mode_hint"] = len(records) < 10

    # 2. 历史清除检测
    history_count = sum(1 for r in records
                        if r.artifact_type == "history")
    cookie_count = sum(1 for r in records
                       if r.artifact_type == "cookie")
    if history_count == 0 and cookie_count > 0:
        indicators["history_cleared"] = True
    # Cookie存在但历史记录为零 → 手动清除

    # 3. 可疑域名关键词匹配
    suspicious_kw = ["torrent","crack","keygen",
                     "warez","pirate","darkweb",
                     "onion","hacktool"]
    for rec in records:
        for kw in suspicious_kw:
            if kw in rec.data.get("url",""):
                suspicious_found.append({...})""")
    add_body_text(s, 7.6, 1.7, 5.2, 5.0, [
        "风险 1: 隐私模式推断",
        "  记录数 < 10 → 可能长期使用无痕模式",
        "  （无痕模式不持久化浏览历史）",
        "",
        "风险 2: 历史清除痕迹",
        "  Cookie 仍在，但 History 为零",
        "  → 用户主动清除了浏览历史",
        "",
        "风险 3: 可疑关键词",
        "  遍历所有 URL 匹配 8 个关键词",
        "  命中后记录 关键词 + URL + 浏览器",
        "  在报告中列出供人工研判",
        "",
        "",
        "以上指标仅作参考，需人工研判",
    ], font_size=12, color=CLR_RED)

    # ── 14. 工具层详解 ──
    s = new_slide(14)
    add_title_bar(s, "工具层详解", "utils/ — 4 个支撑模块")
    add_card(s, 0.5, 1.7, 3.0, 5.2, "sqlite_utils.py", [
        "  safe_connect(path)",
        "  └ mode=ro + query_only=ON",
        "",
        "  safe_connect_with_fallback(path)",
        "  ├ 先尝试只读连接",
        "  └ 失败 → copy_to_temp → 关闭时",
        "     自动 unlink 清理临时文件",
        "",
        "  _dict_factory",
        "  └ 自定义 row_factory 返回 dict",
        "    兼容所有 Python 版本 .get()",
    ])
    add_card(s, 3.8, 1.7, 3.0, 3.0, "time_utils.py", [
        "  chrome_micros_to_iso(µs)",
        "  └ 1601-01-01 epoch 微秒数",
        "    → ISO 8601 字符串",
        "",
        "  firefox_micros_to_iso(µs)",
        "  └ 1970-01-01 epoch (PRTime)",
        "    → ISO 8601 字符串",
        "",
        "  unix_millis_to_iso(ms)",
        "  └ 通用毫秒格式兼容",
    ])
    add_card(s, 7.1, 1.7, 3.0, 2.5, "url_utils.py", [
        "  extract_domain(url)",
        "  └ 从 URL/host 提取纯域名",
        "     去掉协议/端口/前导点",
        "",
        "  classify_url(url)",
        "  └ 匹配 8 个预定义类别",
        "     未命中返回「其他」",
        "",
        "  URL_CATEGORIES dict",
        "  └ 可自定义扩展规则",
    ])
    add_card(s, 10.4, 1.7, 2.7, 2.5, "config.py", [
        "  CHROME_BASE",
        "  EDGE_BASE",
        "  FIREFOX_BASE",
        "  跨平台路径适配",
        "",
        "  OUTPUT_DIR",
        "  HASH_ALGORITHM",
        "  sha256",
    ])

    # ── 15. 图形界面 ──
    s = new_slide(15)
    add_title_bar(s, "图形界面设计", "tkinter 原生 GUI · 零第三方依赖")
    add_body_text(s, 0.7, 1.7, 11.9, 0.8, [
        "布局：标题栏 → 浏览器选择 → 控件面板（分析+输出路径）→ 5 选项卡结果 → 状态栏",
    ], font_size=14, color=CLR_DBLUE)
    add_card(s, 0.5, 2.5, 3.0, 2.0, "控件面板", [
        "  【开始分析】按钮",
        "  └ 启动异步线程分析",
        "",
        "  报告保存位置",
        "  ├ 输入框：默认路径",
        "  └ 【浏览…】目录选择器",
        "",
        "  进度条 + 状态文本",
        "  └ 两阶段进度反馈",
    ])
    add_card(s, 3.8, 2.5, 3.0, 2.0, "5 个选项卡", [
        "  概览      — 记录总数 / 时间范围",
        "             类型分布柱状图",
        "  TOP域名   — 排名表格 Treeview",
        "  活跃时段  — 24h ASCII 热力图",
        "             近 14 天每日统计",
        "  访问类别  — 8 类分类计数",
        "  风险指标  — 三类检测结果",
    ])
    add_card(s, 7.1, 2.5, 5.8, 2.0, "技术要点", [
        "  ▸ threading.Thread 异步执行，UI 不冻结",
        "  ▸ root.after() 线程安全回调更新 UI",
        "  ▸ 分析完成后自动调用 write_report 保存报告",
        "  ▸ 状态栏显示完整保存路径",
        "  ▸ filedialog.askdirectory 选择输出目录",
        "  ▸ python main.py --gui 一键启动",
    ])
    add_code_block(s, 0.7, 5.0, 12, 2.0, """\
thread = threading.Thread(target=self._run_analysis,
                          args=(selected, out_dir), daemon=True)
thread.start()              # 不阻塞主线程

self.root.after(100, lambda: self._analysis_done(True, "分析完成", out_dir))
                            # 线程安全 UI 更新

write_report(self.records, self.profile_result, {}, Path(out_dir))
                            # 自动保存报告到用户指定目录""")

    # ── 16. 报告输出 ──
    s = new_slide(16)
    add_title_bar(s, "报告输出格式", "output/writer.py — 4 类文件")
    add_card(s, 0.5, 1.7, 3.0, 2.5, "forensic_report.json", [
        "  ├─ metadata",
        "  │  ├─ extraction_time",
        "  │  ├─ file_hashes（证据哈希）",
        "  │  └─ total_records",
        "  ├─ profile（7 维度画像）",
        "  └─ records（全部痕迹）",
    ])
    add_card(s, 3.8, 1.7, 3.0, 2.5, "records.csv", [
        "  UTF-8 BOM 编码",
        "  Excel 可直接打开",
        "  ┌────────────────┐",
        "  │ type │url│ts│… │",
        "  │history│…  │…  │…│",
        "  │cookie  │…  │…  │…│",
        "  └────────────────┘",
    ])
    add_card(s, 7.1, 1.7, 3.0, 2.5, "timeline.csv", [
        "  按时间升序排列",
        "  ┌──────────────────┐",
        "  │ 2022-12-20 02:57  │",
        "  │ 2023-03-15 14:22  │",
        "  │ …                 │",
        "  │ 2026-06-17 07:56  │",
        "  └──────────────────┘",
    ])
    add_card(s, 10.4, 1.7, 2.7, 2.5, "report_checksums.json", [
        "  报告自校验文件",
        "  包含 report.json",
        "  records.csv /",
        "  timeline.csv",
        "  各自的 SHA-256",
    ])
    add_body_text(s, 0.7, 4.6, 11.9, 2.5, [
        "报告目录结构示例:",
        '   ~/WebTrail_Forensics/report_20260618_152045/',
        "   ├── forensic_report.json       (完整取证数据)",
        "   ├── records.csv                (扁平化记录表)",
        "   ├── timeline.csv               (时序活动时间线)",
        "   └── report_checksums.json      (自校验哈希)",
    ], font_size=14, color=CLR_DBLUE)

    # ── 17. 实际效果 ──
    s = new_slide(17)
    add_title_bar(s, "实际运行效果", "本地实测数据")
    add_code_block(s, 0.5, 1.7, 6.5, 5.2, """\
==================================================
   WebTrail 数字取证报告摘要
==================================================
  痕迹总数:      1573
  时间线事件:    1570
  时间范围:      2022-12-20 ~ 2026-06-17
  检测浏览器:    Chrome, Edge, Firefox
  配置数:        Default, Profile 1, ...

  TOP 10 域名:
    fanyi.baidu.com               185
    github.com                     93
    file                           82
    accounts.google.com            69
    chatgpt.com                    59
    youtube.com                    51
    cn.bing.com                    39
    www.google.com                 38
    google.com                     36
    claude.ai                      32

  访问类别分布:
    其他               1408
    搜索引擎             77
    技术开发             76
    视频                  9
    社交媒体              2
    邮箱                  1
==================================================""")
    add_card(s, 7.5, 1.7, 5.5, 5.2, "分浏览器统计", [
        "  Chrome :  1379 条记录 · 0 错误",
        "          2 个配置 (Default, Profile 1)",
        "",
        "  Edge :     191 条记录 · 0 错误",
        "          1 个配置 (Default)",
        "",
        "  Firefox :    3 条记录 · 0 错误",
        "          1 个配置 (jwoa6cvq)",
        "",
        "  证据文件:  189 个 SHA-256 哈希",
        "",
        "  画像结论:",
        "  - 活跃高峰: 下午时段（本地时间）",
        "  - 偏好: 技术开发 + 搜索引擎",
        "  - 典型技术人员画像",
    ])

    # ── 18. 尾页 ──
    s = prs.slides.add_slide(BLANK)
    set_slide_bg(s, CLR_COVER_BG)
    add_body_text(s, 1, 1.2, 11, 0.8, ["WebTrail"],
                  font_size=52, color=CLR_WHITE)
    add_body_text(s, 1, 2.0, 11, 0.6, ["项目总结"],
                  font_size=28, color=CLR_LBLUE)
    add_body_text(s, 1, 2.8, 11, 3.5, [
        "核心价值",
        '  「提取 + 分析」一体化取证链路',
        "  从数据采集到报告输出，一个命令完成",
        "",
        "技术特点",
        "  纯 Python 标准库 · 零第三方依赖 · 开箱即用",
        "  15 个源文件 · 约 1800 行代码 · 跨平台",
        "  只读安全 · 链式保管 · 错误隔离",
        "",
        "设计理念",
        "  关注点分离 · 继承复用 · 最少依赖",
        "  取证科学性与工程实用性并重",
        "",
        "扩展性",
        "  新浏览器 → 继承 BaseExtractor",
        "  新分析维度 → 在 profiler 添加函数",
        "  新输出格式 → 在 writer 添加方法",
    ], font_size=15, color=CLR_WHITE)
    add_body_text(s, 1, 6.5, 11, 0.5, [
        "感谢聆听  ·  欢迎提问  ·  https://github.com",
    ], font_size=22, color=CLR_LBLUE)
    add_page_number(s, 18, TOTAL, dark=True)

    # ── 保存 ──
    out = "WebTrail_项目讲解.pptx"
    prs.save(out)
    print(f"PPT 已生成: {out}")


if __name__ == "__main__":
    build()
